import os
import json
import zipfile
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

def emu_to_pixels(emu):
    # Convert EMU (English Metric Unit) to pixels (assuming 96 dpi)
    return int(emu * 96 / 914400)

def extract_media_from_pptx_zip(pptx_path, media_output_folder):
    media_files = []
    try:
        with zipfile.ZipFile(pptx_path, 'r') as pptx_zip:
            media_names = [name for name in pptx_zip.namelist() if name.startswith('ppt/media/')]
            for media_name in media_names:
                media_data = pptx_zip.read(media_name)
                media_filename = os.path.basename(media_name)
                output_path = os.path.join(media_output_folder, media_filename)
                with open(output_path, 'wb') as f:
                    f.write(media_data)
                media_files.append(output_path)
    except Exception as e:
        print(f"Error extracting media from pptx zip: {e}")
    return media_files

def extract_pptx_content(pptx_path, output_folder="extracted_content"):
    prs = Presentation(pptx_path)
    os.makedirs(output_folder, exist_ok=True)
    images_folder = os.path.join(output_folder, "images")
    os.makedirs(images_folder, exist_ok=True)
    media_folder = os.path.join(output_folder, "media")
    os.makedirs(media_folder, exist_ok=True)

    extracted_data = []

    def extract_images_from_shape(shape, slide_index, images_list):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image = shape.image
            image_bytes = image.blob
            image_ext = image.ext
            image_name = f"slide{slide_index+1}_img{len(images_list)+1}.{image_ext}"
            image_path = os.path.join(images_folder, image_name)
            try:
                with open(image_path, "wb") as f:
                    f.write(image_bytes)
                images_list.append(image_path)
                return image_path
            except Exception as e:
                print(f"Error saving image: {e}")
                return None
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            group_images = []
            for shp in shape.shapes:
                img_path = extract_images_from_shape(shp, slide_index, images_list)
                if img_path:
                    group_images.append(img_path)
            return group_images
        return None

    for slide_index, slide in enumerate(prs.slides):
        slide_data = {
            "slide_number": slide_index + 1,
            "title": "",
            "texts": [],
            "images": [],
            "charts": [],
            "tables": [],
            "shapes": []
        }

        # Extract title if present
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            slide_data["title"] = slide.shapes.title.text.strip()

        for shape_index, shape in enumerate(slide.shapes):
            shape_info = {
                "shape_index": shape_index + 1,
                "type": None,
                "text": None,
                "position": {
                    "left": emu_to_pixels(shape.left),
                    "top": emu_to_pixels(shape.top),
                    "width": emu_to_pixels(shape.width),
                    "height": emu_to_pixels(shape.height)
                }
            }

            # Text extraction
            if shape.has_text_frame and shape.text.strip():
                shape_info["type"] = "text"
                shape_info["text"] = shape.text.strip()
                slide_data["texts"].append(shape_info)

            # Image extraction
            image_path = extract_images_from_shape(shape, slide_index, slide_data["images"])
            if image_path:
                shape_info["type"] = "image"
                if isinstance(image_path, list):
                    shape_info["image_paths"] = image_path
                else:
                    shape_info["image_path"] = image_path
                slide_data["images"].append(shape_info)
                continue

            # Chart extraction
            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                shape_info["type"] = "chart"
                # Currently, python-pptx does not support exporting chart images directly
                # So we just note the chart presence and position
                slide_data["charts"].append(shape_info)
                continue

            # Table extraction
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                shape_info["type"] = "table"
                table_data = []
                table = shape.table
                for r in range(len(table.rows)):
                    row_data = []
                    for c in range(len(table.columns)):
                        cell = table.cell(r, c)
                        row_data.append(cell.text.strip())
                    table_data.append(row_data)
                shape_info["table_data"] = table_data
                slide_data["tables"].append(shape_info)
                continue

            # Other shapes
            if shape.shape_type in [MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.GROUP, MSO_SHAPE_TYPE.PLACEHOLDER]:
                shape_info["type"] = "shape"
                # Extract text if any
                if shape.has_text_frame and shape.text.strip():
                    shape_info["text"] = shape.text.strip()
                slide_data["shapes"].append(shape_info)

        extracted_data.append(slide_data)

    # Extract embedded media files from ppt/media folder inside pptx
    media_files = extract_media_from_pptx_zip(pptx_path, media_folder)

    # Add media files info to extracted data summary
    extracted_data.append({
        "embedded_media_files": media_files
    })

    # Save extracted data to JSON file
    json_path = os.path.join(output_folder, "extracted_pptx_content.json")
    with open(json_path, "w", encoding="utf-8") as json_file:
        json.dump(extracted_data, json_file, indent=4, ensure_ascii=False)

    print(f"Extraction complete. Data saved to {json_path}")
    print(f"Extracted images saved to {images_folder}")
    print(f"Extracted media files saved to {media_folder}")

    return json_path, images_folder, media_folder

if __name__ == "__main__":
    import tkinter as tk
    from tkinter import filedialog

    root = tk.Tk()
    root.withdraw()
    pptx_path = filedialog.askopenfilename(
        title="Select a PPTX file",
        filetypes=[("PowerPoint files", "*.pptx")]
    )
    if pptx_path:
        extract_pptx_content(pptx_path)
    else:
        print("No file selected. Exiting.")
