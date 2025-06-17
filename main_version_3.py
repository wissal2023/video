import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from tkinter import Tk, filedialog
from PIL import Image
from io import BytesIO
from pydantic import BaseModel, Field
from typing import List, Optional
import json
from openai import OpenAI
from dotenv import load_dotenv
from gtts import gTTS
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from moviepy import AudioFileClip, ImageClip, concatenate_videoclips
from pdf2image import convert_from_path
import subprocess
import re
import textwrap
from pydantic import BaseModel, ValidationError
from typing import List, Optional,Dict



load_dotenv()

client = OpenAI(
    base_url=os.getenv("ENDPOINT"),
    api_key=os.getenv("TOKEN"),
)
# Data models
class SlideItem(BaseModel):
    title: str
    content: str
    key_points: List[str] = Field(default_factory=list)
    image_prompt: Optional[str] = None

class ShortVideoSegment(BaseModel):
    title: str
    content: str
    script: str
    duration: float = 60.0  # Target duration in seconds

class SlideChunk(BaseModel):
    slides: List[SlideItem]
    voice_over_script: str
    short_segments: List[ShortVideoSegment] = Field(default_factory=list)
    theme_colors: Optional[Dict[str, str]] = None

class VideoConfig(BaseModel):
    theme: str = "professional"  # professional, creative, minimal
    presenter_type: str = "human"  # human, cartoon, none
    language: str = "en"
    voice_style: str = "neutral"  # neutral, enthusiastic, formal
    include_background_music: bool = True
    resolution: str = "1080p"
    aspect_ratio: str = "16:9"
    animation_level: str = "moderate"  # none, subtle, moderate, dynamic




def select_pptx_file():
    root = Tk()
    root.withdraw()
    pptx_path = filedialog.askopenfilename(
        title="Select a PPTX file",
        filetypes=[("PowerPoint files", "*.pptx")]
    )
    return pptx_path

def extract_text_and_images_from_pptx(pptx_path, image_output_folder="pptx_images"):
    print(f"🔍 Extracting text and images from PPTX: {pptx_path}")
    prs = Presentation(pptx_path)
    os.makedirs(image_output_folder, exist_ok=True)

    full_text = ""
    image_paths = []

    for slide_index, slide in enumerate(prs.slides):
        slide_texts = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER or shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                if shape.has_text_frame:
                    text = shape.text.strip()
                    slide_texts.append(text)
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                image_ext = image.ext
                image_name = f"slide{slide_index+1}_img{len(image_paths)+1}.{image_ext}"
                image_path = os.path.join(image_output_folder, image_name)
                with open(image_path, "wb") as f:
                    f.write(image_bytes)
                image_paths.append(image_path)
        full_text += "\n".join(slide_texts) + "\n\n"

    print(f"✅ Extracted text and {len(image_paths)} images from PPTX.")
    return full_text, image_paths


def extract_and_screenshot_shapes(pptx_path, slide_images_folder="slide_images", output_folder="extracted_shapes"):
    """
    Extract images, graphs, shapes, and tables from each slide.
    For pictures, save directly.
    For charts, auto shapes, groups, and tables, crop from slide images based on position and size.
    Save cropped images and metadata with enhanced details.
    """
    from PIL import Image
    import json

    prs = Presentation(pptx_path)
    os.makedirs(output_folder, exist_ok=True)

    # Load slide images
    slide_image_files = sorted([
        os.path.join(slide_images_folder, f) for f in os.listdir(slide_images_folder)
        if f.lower().endswith(('.png', '.jpg', '.jpeg'))
    ])

    metadata = []

    for slide_index, slide in enumerate(prs.slides):
        if slide_index >= len(slide_image_files):
            print(f"Warning: No slide image for slide {slide_index+1}, skipping shape screenshots.")
            continue

        slide_img_path = slide_image_files[slide_index]
        slide_img = Image.open(slide_img_path)

        for shape_index, shape in enumerate(slide.shapes):
            shape_type = shape.shape_type
            shape_name = getattr(shape, "name", f"shape{shape_index+1}")

            # For pictures, save directly
            if shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                image_ext = image.ext
                image_name = f"slide{slide_index+1}_picture{shape_index+1}.{image_ext}"
                image_path = os.path.join(output_folder, image_name)
                with open(image_path, "wb") as f:
                    f.write(image_bytes)
                metadata.append({
                    "slide": slide_index+1,
                    "shape_index": shape_index+1,
                    "type": "picture",
                    "name": shape_name,
                    "image_path": image_path,
                    "position": {
                        "left": shape.left,
                        "top": shape.top,
                        "width": shape.width,
                        "height": shape.height
                    }
                })

            # For charts, auto shapes, groups, and tables, crop from slide image
            elif shape_type in (MSO_SHAPE_TYPE.CHART, MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.GROUP) or shape.has_table:
                # Convert EMU to pixels (assuming 96 dpi)
                def emu_to_px(emu):
                    return int(emu * 96 / 914400)

                left_px = emu_to_px(shape.left)
                top_px = emu_to_px(shape.top)
                width_px = emu_to_px(shape.width)
                height_px = emu_to_px(shape.height)

                # Crop the shape area from slide image
                box = (left_px, top_px, left_px + width_px, top_px + height_px)
                cropped_img = slide_img.crop(box)

                image_name = f"slide{slide_index+1}_{shape_type.name.lower()}{shape_index+1}.png"
                if shape.has_table:
                    image_name = f"slide{slide_index+1}_table{shape_index+1}.png"
                image_path = os.path.join(output_folder, image_name)
                cropped_img.save(image_path)

                style_fill = None
                style_line = None
                # Some shape types like GraphicFrame may not have fill or line attributes
                if hasattr(shape, "fill") and hasattr(shape.fill, "fore_color"):
                    style_fill = shape.fill.fore_color
                if hasattr(shape, "line") and hasattr(shape.line, "color"):
                    style_line = shape.line.color

                metadata.append({
                    "slide": slide_index+1,
                    "shape_index": shape_index+1,
                    "type": "table" if shape.has_table else shape_type.name.lower(),
                    "name": shape_name,
                    "image_path": image_path,
                    "position": {
                        "left": shape.left,
                        "top": shape.top,
                        "width": shape.width,
                        "height": shape.height
                    },
                    "style": {
                        "fill": style_fill,
                        "line": style_line
                    }
                })

    # Save metadata json
    metadata_path = os.path.join(output_folder, "extracted_shapes_metadata.json")
    with open(metadata_path, "w") as f:
        json.dump(metadata, f, indent=4)

    print(f"✅ Extracted and saved shapes, tables, and images to {output_folder}")
    print(f"Metadata saved to {metadata_path}")
    return metadata, output_folder

def extract_json_from_response(response: str) -> str:
    # Remove introductory text before the first JSON-like bracket
    start = response.find('{')
    end = response.rfind('}') + 1
    if start != -1 and end != -1:
        return response[start:end]
    return response  # fallback (will fail as-is)


def chunk_text(text, chunk_size=100000):
    print("Chunking text...")
    chunks = textwrap.wrap(text, chunk_size, break_long_words=False)
    print(f"✅ Created {len(chunks)} chunks.")
    return chunks

def generate_chunk_content(chunks, config):
    print("Summarizing chunks with Mistral...")
    combined_content = "\n\n".join(chunks)

    theme_desc = {
        "professional": "formal, corporate style with clean design",
        "creative": "vibrant, engaging style with dynamic elements",
        "minimal": "clean, simple style with focus on key content"
    }.get(config.theme, "professional style")

    voice_desc = {
        "neutral": "balanced and clear",
        "enthusiastic": "energetic and engaging",
        "formal": "serious and professional"
    }.get(config.voice_style, "clear and professional")

    prompt = (
        f"Generate a structured presentation based on the following content. "
        f"Use a {theme_desc} visual approach and a {voice_desc} tone for narration.\n\n"
        "Create a JSON with these keys:\n"
        "1. 'slides': list of objects with 'title', 'content', 'key_points' (list of bullet points), "
        "and 'image_prompt' (a description for generating a relevant image)\n"
        "2. 'voice_over_script': professional narration script covering all slides\n"
        "3. 'short_segments': 3-5 stand-alone segments for short-form videos (under 2 minutes each) "
        "with 'title', 'content', 'script', and 'duration' fields\n"
        "4. 'theme_colors': suggested color scheme (primary, secondary, accent, background, text)\n\n"
        f"Content:\n{combined_content}\n\n"
        "Respond with valid JSON only. Keep all content factual and based on the input material. "
        "Ensure 'voice_over_script' is a single string, not a list."
    )

    response = client.chat.completions.create(
        model="mistral-medium",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.3,
    ).choices[0].message.content.strip()

    # Clean up response to get valid JSON
    if response.startswith("```json"):
        response = response.lstrip("```json").rstrip("```").strip()
    elif response.startswith("```"):
        response = response.lstrip("```").rstrip("```").strip()

    try:
        parsed_response = json.loads(response)

        # ✅ Join voice_over_script if it is a list
        if isinstance(parsed_response.get('voice_over_script'), list):
            parsed_response['voice_over_script'] = "\n\n".join(parsed_response['voice_over_script'])

        # ✅ Convert complex duration strings into float seconds
        
        def parse_duration_to_seconds(text):
            if isinstance(text, (int, float)):
                return float(text)
            if isinstance(text, str):
                match = re.match(r"(?:(\d+)\s*minutes?)?\s*(?:(\d+)\s*seconds?)?", text.lower())
                if match:
                    minutes = int(match.group(1)) if match.group(1) else 0
                    seconds = int(match.group(2)) if match.group(2) else 0
                    return float(minutes * 60 + seconds)
            return 60.0  # default fallback

        for seg in parsed_response.get("short_segments", []):
            seg["duration"] = parse_duration_to_seconds(seg.get("duration"))

        validated_chunk = SlideChunk(**parsed_response)
        print(validated_chunk)

    except (json.JSONDecodeError, ValidationError) as e:
        print(f"Parsing error: {e}\nResponse was: {response}")
        raise

    print("✅ Summarization complete.")
    return validated_chunk

# Step 3: Generate Audio from Script
def generate_audio(script, output_file):
    tts = gTTS(script, lang="fr")
    tts.save(output_file)


def generate_presentation(slide_contents, pptx_path, config=None, slide_images=None):
    """
    Generate a PowerPoint presentation from slide contents, which can be:
    1. A list of SlideItem objects
    2. A SlideChunk object
    3. The raw parsed content from the paste.txt
    
    Parameters:
    - slide_contents: Slide content in one of the above formats
    - pptx_path: Path where the presentation will be saved
    - config: Optional configuration parameters
    
    Returns:
    - Path to the saved presentation
    """

    
    print("Creating enhanced slides...")
    prs = Presentation()
    
    # Handle different possible input formats
    if hasattr(slide_contents, 'slides'):
        # This is a SlideChunk object
        slides = slide_contents.slides
        voice_over_script = slide_contents.voice_over_script
        short_segments = slide_contents.short_segments
        theme_colors = slide_contents.theme_colors if slide_contents.theme_colors else {}
    elif isinstance(slide_contents, list) and all(hasattr(item, 'title') for item in slide_contents):
        # This is a list of SlideItem objects
        slides = slide_contents
        voice_over_script = ""
        short_segments = []
        theme_colors = {}
    elif isinstance(slide_contents, list) and len(slide_contents) > 0 and 'slides' in dir(slide_contents[0]):
        # This might be a list containing one SlideChunk object
        slides = slide_contents[0].slides
        voice_over_script = slide_contents[0].voice_over_script
        short_segments = slide_contents[0].short_segments
        theme_colors = slide_contents[0].theme_colors if slide_contents[0].theme_colors else {}
    else:
        # Assume raw data format like in paste.txt
        # Try to extract slides and other info
        slides = slide_contents
        voice_over_script = ""
        short_segments = []
        theme_colors = {}
        
        # Check if it's the raw parsed structure
        if hasattr(slide_contents, 'voice_over_script'):
            voice_over_script = slide_contents.voice_over_script
        if hasattr(slide_contents, 'short_segments'):
            short_segments = slide_contents.short_segments
        if hasattr(slide_contents, 'theme_colors'):
            theme_colors = slide_contents.theme_colors
    
    # Ensure we have default theme colors if not provided
    if not theme_colors:
        theme_colors = {
            "primary": "#1F497D",
            "secondary": "#4F81BD", 
            "accent": "#C0504D",
            "background": "#FFFFFF",
            "text": "#000000"
        }
    
    # Strip '#' prefix from colors if present
    for key in theme_colors:
        if isinstance(theme_colors[key], str) and theme_colors[key].startswith('#'):
            theme_colors[key] = theme_colors[key][1:]
    
    # Add a title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    
    # Determine presentation title from first slide
    presentation_title = "Presentation"
    if slides and hasattr(slides[0], 'title'):
        first_title = slides[0].title
        if "Introduction" in first_title and "to" in first_title:
            presentation_title = first_title.split("to")[1].strip()
        else:
            presentation_title = first_title
    
    title.text = presentation_title
    subtitle.text = "A Comprehensive Guide"
    
    # Style the title slide
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(theme_colors["primary"])
    subtitle.text_frame.paragraphs[0].font.size = Pt(28)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(theme_colors["secondary"])
    
    # Process each slide
    for i, slide_item in enumerate(slides):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_shape = slide.shapes.title
        content_placeholder = slide.placeholders[1]
        
        # Get slide title
        slide_title = f"Slide {i+1}"
        if hasattr(slide_item, 'title'):
            slide_title = slide_item.title
        
        # Set the title
        title_shape.text = slide_title
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(theme_colors["primary"])
        
        # Set the content
        content_frame = content_placeholder.text_frame
        content_frame.clear()
        
        # Add main content
        slide_content = "Content"
        if hasattr(slide_item, 'content'):
            slide_content = slide_item.content
            
        p = content_frame.add_paragraph()
        p.text = slide_content
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor.from_string(theme_colors["text"])
        
        # Add key points as bullets
        key_points = []
        if hasattr(slide_item, 'key_points'):
            key_points = slide_item.key_points
            
        if key_points:
            content_frame.add_paragraph().text = ""  # Add spacing
            for point in key_points:
                bullet_p = content_frame.add_paragraph()
                bullet_p.text = point
                bullet_p.font.size = Pt(20)
                bullet_p.level = 1
                bullet_p.font.color.rgb = RGBColor.from_string(theme_colors["secondary"])
        
        # Add images to slide if available
        if slide_images and i < len(slide_images):
            img_path = slide_images[i]
            if isinstance(img_path, str):
                left = Inches(5.5)
                top = Inches(1.5)
                height = Inches(3.5)
                slide.shapes.add_picture(img_path, left, top, height=height)
            else:
                print(f"Warning: slide_images[{i}] is not a file path string, skipping image.")
    
   
    
    # Add a final slide
    final_slide = prs.slides.add_slide(prs.slide_layouts[2])
    final_title = final_slide.shapes.title
    
    final_title.text = "you can now go to quiz"
    final_title.text_frame.paragraphs[0].font.size = Pt(40)
    final_title.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(theme_colors["primary"])
    
    
    
    # Save the presentation
    prs.save(pptx_path)
    print(f"✅ Enhanced slides created and saved to {pptx_path}")
    return pptx_path

# Step 5: Convert Slides to Images
def slides_to_images(ppt_path, output_folder):
    libreoffice_path = r"C:\Program Files\LibreOffice\program\soffice.exe"  
    subprocess.run([
        libreoffice_path, '--headless', '--convert-to', 'pdf', ppt_path, '--outdir', output_folder
    ], check=True)

    pdf_path = os.path.join(output_folder, os.path.splitext(os.path.basename(ppt_path))[0] + ".pdf")

    return [img.save(os.path.join(output_folder, f"slide_{i}.png"), 'PNG') or os.path.join(output_folder, f"slide_{i}.png")
            for i, img in enumerate(convert_from_path(pdf_path, dpi=200))]

# Step 6: Generate Video

def create_video(slide_imgs, audio_path, output_path):
    audio = AudioFileClip(audio_path)
    duration = audio.duration / len(slide_imgs)
    clips = [ImageClip(p).with_duration(duration) for p in slide_imgs]
    video = concatenate_videoclips(clips, method="compose").with_audio(audio)
    video.write_videofile(output_path, fps=24)


def summarize_with_mistral(text, chunks, config, slide_images=None):

    slide_chunk = generate_chunk_content(chunks, config)

    content = "\n\n".join([f"Title: {s.title}\nContent: {s.content}" for s in slide_chunk.slides])
    theme_desc = {
        "professional": "formal, corporate style with clean design",
        "creative": "vibrant, engaging style with dynamic elements",
        "minimal": "clean, simple style with focus on key content"
    }.get(config.theme, "professional style")

    voice_desc = {
        "neutral": "balanced and clear",
        "enthusiastic": "energetic and engaging",
        "formal": "serious and professional"
    }.get(config.voice_style, "clear and professional")
    combined_content = "\n\n".join(chunks)

    prompt = (
        f"Generate a structured presentation based on the following content. "
        f"Use a {theme_desc} visual approach and a {voice_desc} tone for narration.\n\n"
        "Create a JSON with these keys:\n"
        "1. 'slides': list of objects with 'title', 'content', 'key_points' (list of bullet points), "
        "and 'image_prompt' (a description for generating a relevant image)\n"
        "2. 'voice_over_script': professional narration script covering all slides\n"
        "3. 'short_segments': 3-5 stand-alone segments for short-form videos (under 2 minutes each) "
        "with 'title', 'content', 'script', and 'duration' fields\n"
        "4. 'theme_colors': suggested color scheme (primary, secondary, accent, background, text)\n\n"
        f"Content:\n{combined_content}\n\n"
        "Respond with valid JSON only. Keep all content factual and based on the input material. "
        "Ensure 'voice_over_script' is a single string, not a list."
    )

    print("Sending request to Mistral API...")
    try:
        response = client.chat.completions.create(
            model="mistral-medium",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.3,
        ).choices[0].message.content.strip()
        print("Received response from Mistral API.")
    except Exception as e:
        print(f"Error calling Mistral API: {e}")
        raise

    print("\n✅ Mistral Summary:\n")
    print(response)

    # Parse JSON response
    cleaned_json = extract_json_from_response(response)
    try:
        data = json.loads(cleaned_json)
    except json.JSONDecodeError as e:
        print("❌ Failed to parse cleaned JSON:")
        print(cleaned_json[:1000])  # print preview of what failed
        print(f"Error: {e}")
        return


    # Generate audio narration from voice_over_script
    voice_script = data.get("voice_over_script", "")
    if voice_script:
        print("Generating audio narration...")
        # voice_script might be a list, join if so
        if isinstance(voice_script, list):
            voice_script = "\n\n".join(voice_script)
        elif not isinstance(voice_script, str):
            voice_script = str(voice_script)
        tts = gTTS(voice_script, lang="fr")
        audio_file = "voice_narration.mp3"
        tts.save(audio_file)
        print(f"Audio narration saved to {audio_file}")
    else:
        print("No voice_over_script found, skipping audio generation.")
        audio_file = None

    # Generate PowerPoint presentation
    print("Generating PowerPoint presentation...")
    prs = Presentation()
    theme_colors = data.get("theme_colors", {})
    slides = data.get("slides", [])

    # Default colors if not provided
    default_colors = {
        "primary": "003366",
        "secondary": "0066cc",
        "background": "f8f9fa",
        "text": "333333",
        "accent": "ff6600"
    }
    for key in default_colors:
        if key not in theme_colors or not theme_colors[key]:
            theme_colors[key] = default_colors[key]

    # Remove leading '#' if present in color codes
    for key in theme_colors:
        if isinstance(theme_colors[key], str) and theme_colors[key].startswith('#'):
            theme_colors[key] = theme_colors[key][1:]

    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    title.text = slides[0]["title"] if slides else "Presentation"
    subtitle.text = "Generated Presentation"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(theme_colors["primary"])
    subtitle.text_frame.paragraphs[0].font.size = Pt(28)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(theme_colors["secondary"])

    # Add slides
    for i, slide_data in enumerate(slides):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_shape = slide.shapes.title
        content_placeholder = slide.placeholders[1]

        title_shape.text = slide_data.get("title", "Slide")
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(theme_colors["primary"])

        content_frame = content_placeholder.text_frame
        content_frame.clear()

        content_text = slide_data.get("content", "")
        if isinstance(content_text, list):
            # content_text list may contain dicts, convert all to strings
            content_text = "\n".join(str(item) if isinstance(item, dict) else item for item in content_text)
        p = content_frame.add_paragraph()
        p.text = content_text
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor.from_string(theme_colors["text"])

        key_points = slide_data.get("key_points", [])
        if key_points:
            content_frame.add_paragraph().text = ""
            for point in key_points:
                bullet_p = content_frame.add_paragraph()
                bullet_p.text = point
                bullet_p.font.size = Pt(20)
                bullet_p.level = 1
                bullet_p.font.color.rgb = RGBColor.from_string(theme_colors["secondary"])

    pptx_file = "generated_presentation.pptx"
    prs.save(pptx_file)
    print(f"PowerPoint presentation saved to {pptx_file}")

    # Convert slides to images using LibreOffice
    print("Converting slides to images...")
    output_folder = "slide_images"
    os.makedirs(output_folder, exist_ok=True)
    libreoffice_path = r"C:\Program Files\LibreOffice\program\soffice.exe"
    try:
        subprocess.run([
            libreoffice_path, '--headless', '--convert-to', 'pdf', pptx_file, '--outdir', output_folder
        ], check=True)
    except Exception as e:
        print(f"Error converting PPTX to PDF: {e}")
        return

    pdf_path = os.path.join(output_folder, os.path.splitext(os.path.basename(pptx_file))[0] + ".pdf")

    try:
        images = convert_from_path(pdf_path, dpi=200)
    except Exception as e:
        print(f"Error converting PDF to images: {e}")
        return

    slide_image_paths = []
    for i, img in enumerate(images):
        img_path = os.path.join(output_folder, f"slide_{i+1}.png")
        img.save(img_path, "PNG")
        slide_image_paths.append(img_path)

    # Create video from images and audio
    if audio_file:
        print("Creating video from slides and audio...")
        audio_clip = AudioFileClip(audio_file)
        duration_per_slide = audio_clip.duration / len(slide_image_paths)
        clips = [ImageClip(img).with_duration(duration_per_slide) for img in slide_image_paths]
        video = concatenate_videoclips(clips, method="compose").with_audio(audio_clip)
        video_file = "final_video.mp4"
        video.write_videofile(video_file, fps=24)
        print(f"Video saved to {video_file}")
    else:
        print("No audio file available, skipping video creation.")


if __name__ == "__main__":
    pptx_file = select_pptx_file()
    if pptx_file:
        slide_items = extract_text_and_images_from_pptx(pptx_file)
        # extract_text_and_images_from_pptx returns (text, images), but summarize_with_mistral expects slide_items list
        # So we need to adapt the call accordingly
        text, images = slide_items
        chunks = chunk_text(text)
        config = VideoConfig()  # Use default config or customize as needed
        summarize_with_mistral(text, chunks, config, images)

        # Convert slides to images first
        slide_images_folder = "slide_images"
        os.makedirs(slide_images_folder, exist_ok=True)
        slides_to_images(pptx_file, slide_images_folder)

        # Extract and screenshot shapes and graphs
        extract_and_screenshot_shapes(pptx_file, slide_images_folder, output_folder="extracted_shapes")

    else:
        print("❌ No file selected.")
